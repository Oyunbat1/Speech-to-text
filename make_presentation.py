"""Generate the Mongolian presentation as a .pptx file."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

ACCENT = RGBColor(0x25, 0x63, 0xEB)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
MUTED = RGBColor(0x6B, 0x72, 0x80)
BG = RGBColor(0xF7, 0xF7, 0xF8)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]


def add_bg(slide):
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    return bg


def add_text(slide, left, top, width, height, text, *,
             size=18, bold=False, color=DARK, align=None):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if align is not None:
            p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = "Calibri"
    return tb


def add_accent_bar(slide, top=Inches(1.1)):
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.7), top, Inches(0.5), Inches(0.08)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()


def bullet_slide(title, bullets, *, subtitle=None):
    slide = prs.slides.add_slide(BLANK)
    add_bg(slide)
    add_text(slide, Inches(0.7), Inches(0.5), Inches(12), Inches(0.7),
             title, size=32, bold=True, color=DARK)
    add_accent_bar(slide)
    if subtitle:
        add_text(slide, Inches(0.7), Inches(1.25), Inches(12), Inches(0.5),
                 subtitle, size=16, color=MUTED)
    top = Inches(2.0) if subtitle else Inches(1.6)
    tb = slide.shapes.add_textbox(Inches(0.9), top, Inches(11.5), Inches(5))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(10)
        run = p.add_run()
        run.text = f"•  {b}"
        run.font.size = Pt(20)
        run.font.color.rgb = DARK
        run.font.name = "Calibri"


# Slide 1 — Title
slide = prs.slides.add_slide(BLANK)
add_bg(slide)
band = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, 0, Inches(2.8), prs.slide_width, Inches(2.0)
)
band.fill.solid()
band.fill.fore_color.rgb = ACCENT
band.line.fill.background()
add_text(slide, Inches(0.7), Inches(3.0), Inches(12), Inches(1.0),
         "Монгол хэлний хичээл хөрвүүлэгч",
         size=44, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
add_text(slide, Inches(0.7), Inches(4.0), Inches(12), Inches(0.6),
         "Хичээлийн видеог автоматаар Excel болгон хөрвүүлэх вэб апп",
         size=20, color=RGBColor(0xE0, 0xE7, 0xFF))
add_text(slide, Inches(0.7), Inches(6.5), Inches(12), Inches(0.4),
         "MVP танилцуулга  •  2026",
         size=14, color=MUTED)

# Slide 2 — Problem
bullet_slide(
    "Шийдэж буй асуудал",
    [
        "Багш нар хичээлээ бичлэг болгож хадгалдаг ч буцаж сонсох цаг шаардсан.",
        "Хичээлийн агуулгыг бичгээр гаргахад гараар тэмдэглэл хийх нь удаан.",
        "Багшийн тайлбар ба сурагчийн хариултыг ялгах нь нэмэлт ажил.",
        "Шийдэл: автомат хөрвүүлэлт + яригч таних → бэлэн Excel файл.",
    ],
    subtitle="Яагаад энэ төсөл хэрэгтэй вэ?",
)

# Slide 3 — How it works
bullet_slide(
    "Хэрхэн ажилладаг вэ?",
    [
        "1. Багшийн дуу хоолойн жишээ (10-30 сек) оруулна.",
        "2. Хичээлийн видео/аудио файлаа оруулна.",
        "3. ffmpeg → 16kHz моно WAV формат руу хөрвүүлэх.",
        "4. pyannote.audio → яригч салгалт (diarization).",
        "5. Дуу хоолойн эмбеддинг + косинусын төстэй байдал → Багш таних.",
        "6. Chimege API → монгол хэл рүү бичвэр болгох.",
        "7. openpyxl → эцсийн Excel файл: Start | End | Role | Text.",
    ],
    subtitle="Хоёр файл оруулаад, нэг Excel хүлээн авна",
)

# Slide 4 — Tech stack
bullet_slide(
    "Технологийн стек",
    [
        "Python 3.10+  •  FastAPI  •  Uvicorn",
        "pyannote.audio 3.x — яригч салгалт ба эмбеддинг",
        "Chimege API — монгол хэлний speech-to-text",
        "ffmpeg + pydub — аудио боловсруулалт",
        "openpyxl — Excel файл бичих",
        "Энгийн HTML + Vanilla JS — UI (framework байхгүй)",
    ],
    subtitle="Хөнгөн, цэвэр, MVP-д тохирсон",
)

# Slide 5 — Demo
slide = prs.slides.add_slide(BLANK)
add_bg(slide)
add_text(slide, Inches(0.7), Inches(0.5), Inches(12), Inches(0.7),
         "Демо", size=32, bold=True, color=DARK)
add_accent_bar(slide)
add_text(slide, Inches(0.7), Inches(1.25), Inches(12), Inches(0.5),
         "Амьд үзүүлэлт — localhost:8000", size=16, color=MUTED)

# Demo steps as numbered cards
steps = [
    ("1", "Файл оруулах",
     "Багшийн жишээ + хичээлийн бичлэг"),
    ("2", "Боловсруулах",
     "Audio → Diarize → Match → STT"),
    ("3", "Excel татах",
     "Start | End | Role | Text"),
]
card_w = Inches(3.8)
card_h = Inches(3.2)
gap = Inches(0.3)
total_w = card_w * 3 + gap * 2
start_left = (prs.slide_width - total_w) / 2
for i, (num, title, desc) in enumerate(steps):
    left = start_left + (card_w + gap) * i
    top = Inches(2.5)
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, card_w, card_h)
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    card.line.color.rgb = RGBColor(0xE5, 0xE7, 0xEB)
    add_text(slide, left, top + Inches(0.3), card_w, Inches(0.8),
             num, size=48, bold=True, color=ACCENT, align=2)
    add_text(slide, left, top + Inches(1.4), card_w, Inches(0.5),
             title, size=20, bold=True, color=DARK, align=2)
    add_text(slide, left + Inches(0.3), top + Inches(2.1), card_w - Inches(0.6), Inches(1.0),
             desc, size=14, color=MUTED, align=2)

# Slide 6 — Limitations & future
bullet_slide(
    "Хязгаарлалт ба дараагийн алхам",
    [
        "Зөвхөн монгол хэл (хэл хольсон яриа алдаатай).",
        "Багшийн жишээ цэвэр, чимээгүй байх шаардлагатай.",
        "CPU дээр 60 мин хичээл ≈ 50 мин боловсруулах хугацаа.",
        "Цаашид: GPU дэмжлэг, бодит цагийн прогресс, олон багштай горим.",
    ],
    subtitle="Юу ажиллахгүй байна, юу нэмж болох вэ",
)

# Slide 7 — Thank you
slide = prs.slides.add_slide(BLANK)
add_bg(slide)
band = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, 0, Inches(2.8), prs.slide_width, Inches(2.0)
)
band.fill.solid()
band.fill.fore_color.rgb = ACCENT
band.line.fill.background()
add_text(slide, Inches(0.7), Inches(3.1), Inches(12), Inches(1.0),
         "Баярлалаа",
         size=54, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
add_text(slide, Inches(0.7), Inches(4.1), Inches(12), Inches(0.6),
         "Асуулт байвал хүлээж авна",
         size=22, color=RGBColor(0xE0, 0xE7, 0xFF))

out = "D:/speech-to-text/presentation.pptx"
prs.save(out)
print(f"Saved: {out}")
